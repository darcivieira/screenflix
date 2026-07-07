#!/usr/bin/env python3
"""Gera o deck do Workshop ScreenFlix (.pptx) — identidade visual ScreenFlix."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Paleta ScreenFlix ----
INK      = RGBColor(0x0B, 0x0B, 0x0F)   # fundo escuro
SURF1    = RGBColor(0x14, 0x14, 0x1D)
SURF2    = RGBColor(0x1E, 0x1E, 0x2A)
BRAND    = RGBColor(0xE5, 0x09, 0x14)   # vermelho
BRAND_H  = RGBColor(0xF6, 0x12, 0x1D)
ACCENT   = RGBColor(0x00, 0xB3, 0xFF)
TXT      = RGBColor(0xF5, 0xF7, 0xFA)
MUTED    = RGBColor(0xA9, 0xB0, 0xBD)
OK       = RGBColor(0x22, 0xC5, 0x5E)
WARN     = RGBColor(0xF5, 0x9E, 0x0B)
CODE_BG  = RGBColor(0x0F, 0x0F, 0x17)
CARD     = RGBColor(0x1A, 0x1A, 0x25)

FONT   = "Arial"
MONO   = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def bg(slide, color=INK):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _set_fill(r, color)
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def rect(slide, x, y, w, h, color, line=None, line_w=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    if line is not None:
        r.line.color.rgb = line; r.line.width = line_w or Pt(1)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    return r


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (text, size, color, bold, font)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, fnt) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = fnt or FONT
    return tb


def R(t, sz=18, col=TXT, bold=False, fnt=None):
    return (t, sz, col, bold, fnt)


def kicker_bar(slide, x, y, h=Inches(0.55), color=BRAND):
    rect(slide, x, y, Inches(0.09), h, color)


def code_box(slide, x, y, w, h, lines, size=13):
    box = rect(slide, x, y, w, h, CODE_BG, line=SURF2, line_w=Pt(1))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (line, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.12; p.space_after = Pt(0)
        r = p.add_run(); r.text = line if line else " "
        r.font.name = MONO; r.font.size = Pt(size); r.font.color.rgb = col
    return box


def footer(slide, day_label):
    txt(slide, Inches(0.55), Inches(7.02), Inches(8), Inches(0.35),
        [[R("ScreenFlix Workshop", 10, MUTED, True), R("   ·   " + day_label, 10, MUTED)]])
    txt(slide, Inches(11.3), Inches(7.02), Inches(1.5), Inches(0.35),
        [[R("", 10, MUTED)]], align=PP_ALIGN.RIGHT)


def page_num(slide, n):
    txt(slide, Inches(12.5), Inches(7.02), Inches(0.6), Inches(0.35),
        [[R(str(n), 10, MUTED)]], align=PP_ALIGN.RIGHT)


_counter = {"n": 0}

def content_slide(day_label, title, subtitle=None, accent=BRAND):
    _counter["n"] += 1
    s = prs.slides.add_slide(BLANK)
    bg(s, INK)
    kicker_bar(s, Inches(0.55), Inches(0.55), Inches(0.5), accent)
    txt(s, Inches(0.78), Inches(0.5), Inches(11.8), Inches(0.4),
        [[R(day_label.upper(), 13, accent, True)]])
    txt(s, Inches(0.78), Inches(0.86), Inches(11.9), Inches(0.9),
        [[R(title, 30, TXT, True)]])
    if subtitle:
        txt(s, Inches(0.78), Inches(1.62), Inches(11.9), Inches(0.5),
            [[R(subtitle, 15, MUTED)]])
    footer(s, day_label)
    page_num(s, _counter["n"])
    return s


def bullets(slide, x, y, w, items, size=17, gap=10, col=TXT, marker_col=BRAND):
    tb = slide.shapes.add_textbox(x, y, w, Inches(4))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.08
        # marker
        m = p.add_run(); m.text = "▸  "
        m.font.name = FONT; m.font.size = Pt(size); m.font.color.rgb = marker_col; m.font.bold = True
        if isinstance(item, tuple):
            head, rest = item
            r = p.add_run(); r.text = head
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = TXT; r.font.bold = True
            if rest:
                r2 = p.add_run(); r2.text = rest
                r2.font.name = FONT; r2.font.size = Pt(size); r2.font.color.rgb = MUTED
        else:
            r = p.add_run(); r.text = item
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = col
    return tb


def card(slide, x, y, w, h, title, body_lines, accent=BRAND, title_size=16, body_size=13):
    c = rect(slide, x, y, w, h, CARD, line=SURF2, line_w=Pt(1))
    rect(slide, x, y, w, Inches(0.07), accent)
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]; p.space_after = Pt(6)
    r = p.add_run(); r.text = title
    r.font.name = FONT; r.font.size = Pt(title_size); r.font.bold = True; r.font.color.rgb = TXT
    for line in body_lines:
        p = tf.add_paragraph(); p.space_after = Pt(3); p.line_spacing = 1.05
        r = p.add_run(); r.text = line
        r.font.name = FONT; r.font.size = Pt(body_size); r.font.color.rgb = MUTED
    return c


def pill(slide, x, y, text, color=BRAND, w=Inches(1.6)):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.42))
    p.fill.solid(); p.fill.fore_color.rgb = color; p.line.fill.background()
    p.shadow.inherit = False
    try:
        p.adjustments[0] = 0.5
    except Exception:
        pass
    tf = p.text_frame; tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    r = para.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = TXT
    return p


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, INK)
rect(s, 0, 0, SW, Inches(0.16), BRAND)
rect(s, 0, Inches(7.34), SW, Inches(0.16), BRAND)
txt(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.6),
    [[R("S C R E E N F L I X", 22, BRAND, True)]])
txt(s, Inches(0.85), Inches(2.15), Inches(11.8), Inches(1.6),
    [[R("Workshop de Engenharia", 46, TXT, True)]])
txt(s, Inches(0.9), Inches(3.35), Inches(11.6), Inches(0.6),
    [[R("Do container ao código, do código à IA — um projeto como fio condutor.", 18, MUTED)]])
# 3 day cards
cy = Inches(4.35); cw = Inches(3.75); ch = Inches(1.5); gap = Inches(0.28)
cx0 = Inches(0.9)
days = [("DIA 1", "Infra & Containers", BRAND),
        ("DIA 2", "Programação Pragmática", ACCENT),
        ("DIA 3", "IA no Desenvolvimento", OK)]
for i, (d, t, ac) in enumerate(days):
    x = cx0 + i * (cw + gap)
    card(s, x, cy, cw, ch, "", [], accent=ac)
    txt(s, x + Inches(0.22), cy + Inches(0.28), cw - Inches(0.4), Inches(0.4),
        [[R(d, 13, ac, True)]])
    txt(s, x + Inches(0.22), cy + Inches(0.62), cw - Inches(0.4), Inches(0.7),
        [[R(t, 19, TXT, True)]])
txt(s, Inches(0.9), Inches(6.35), Inches(11.6), Inches(0.6),
    [[R("3 encontros de 1h30", 13, MUTED, True),
      R("     ·     FastAPI · PostgreSQL · Docker · Vue 3 · Claude Code", 13, MUTED)]])

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════
s = content_slide("Visão Geral", "Agenda dos 3 dias")
data = [
    ("DIA 1 · Infra & Containers", BRAND,
     ["Por que containers", "Dockerfile multi-stage", "docker-compose", "Comandos & Makefile"]),
    ("DIA 2 · Programação Pragmática", ACCENT,
     ["Ambientes modulares", "Clean Architecture", "Regra de dependência", "Trade-offs pragmáticos"]),
    ("DIA 3 · IA no Desenvolvimento", OK,
     ["Claude Code & CLAUDE.md", "Skills", "Subagents", "Agent Teams"]),
]
cw = Inches(3.95); ch = Inches(3.9); gap = Inches(0.3); x0 = Inches(0.6); y0 = Inches(2.3)
for i, (title, ac, items) in enumerate(data):
    x = x0 + i * (cw + gap)
    card(s, x, y0, cw, ch, title, [], accent=ac, title_size=16)
    bullets(s, x + Inches(0.22), y0 + Inches(0.95), cw - Inches(0.45), items,
            size=15, gap=13, marker_col=ac)
txt(s, Inches(0.6), Inches(6.35), Inches(12), Inches(0.5),
    [[R("Um único código-base — o ScreenFlix — conecta infra, arquitetura e IA.", 14, MUTED)]])

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — O PROJETO
# ══════════════════════════════════════════════════════════════
s = content_slide("Introdução", "O projeto ScreenFlix", "Catálogo de filmes e séries — nosso estudo de caso vivo")
rows = [
    ("Backend", "FastAPI · SQLAlchemy async · Pydantic · Structlog · uvicorn"),
    ("Dados", "PostgreSQL 18 — tabelas media e episode"),
    ("Infra", "Docker (multi-stage) · docker-compose · uv"),
    ("Frontend", "Vue 3 · Vue Router · Pinia · Vite"),
    ("Integrações", "OMDb (ingestão) + OpenAI (normalização por schema)"),
]
y = Inches(2.35)
for i, (k, v) in enumerate(rows):
    yy = y + i * Inches(0.82)
    rect(s, Inches(0.6), yy, Inches(2.6), Inches(0.66), SURF2)
    txt(s, Inches(0.78), yy + Inches(0.16), Inches(2.4), Inches(0.4), [[R(k, 15, ACCENT, True)]])
    rect(s, Inches(3.3), yy, Inches(9.4), Inches(0.66), CARD)
    txt(s, Inches(3.5), yy + Inches(0.16), Inches(9.1), Inches(0.4), [[R(v, 14, TXT)]])

# ══════════════════════════════════════════════════════════════
# DIA 1
# ══════════════════════════════════════════════════════════════
# Divider
def divider(day, title, items, accent):
    s = prs.slides.add_slide(BLANK)
    bg(s, INK)
    rect(s, 0, 0, Inches(0.35), SH, accent)
    txt(s, Inches(0.95), Inches(2.1), Inches(11), Inches(0.6), [[R(day.upper(), 22, accent, True)]])
    txt(s, Inches(0.9), Inches(2.75), Inches(11.6), Inches(1.2), [[R(title, 42, TXT, True)]])
    bullets(s, Inches(0.98), Inches(4.35), Inches(11), items, size=17, gap=9, marker_col=accent)
    return s

divider("Dia 1", "Infra & Containers",
        ["Empacotar a aplicação de forma reprodutível",
         "Dockerfile multi-stage · docker-compose · comandos · Makefile"], BRAND)

# 1.1 por que containers
s = content_slide("Dia 1 · Infra & Containers", "Por que containers?")
txt(s, Inches(0.6), Inches(2.15), Inches(12), Inches(0.6),
    [[R("Sem containers, cada dev reproduz o ambiente na mão — a origem do ", 16, MUTED),
      R("“na minha máquina funciona”.", 16, TXT, True)]])
cw = Inches(3.95); ch = Inches(2.5); gap = Inches(0.3); x0 = Inches(0.6); y0 = Inches(3.0)
cards = [("Reprodutibilidade", ["A mesma imagem roda igual", "em qualquer lugar."]),
         ("Isolamento", ["O Postgres do projeto não", "conflita com outro local."]),
         ("Paridade dev/prod", ["A imagem testada é a", "que vai para produção."])]
for i, (t, b) in enumerate(cards):
    card(s, x0 + i*(cw+gap), y0, cw, ch, t, b, accent=BRAND, body_size=14)

# 1.2 Dockerfile
s = content_slide("Dia 1 · Infra & Containers", "Dockerfile multi-stage")
code_box(s, Inches(0.6), Inches(2.15), Inches(7.1), Inches(4.4), [
    ("# Stage 1: builder", MUTED),
    ("FROM python:3.13-slim AS builder", TXT),
    ("RUN apt-get install build-essential libpq-dev", TXT),
    ("COPY pyproject.toml uv.lock ./   # cache de deps", ACCENT),
    ("RUN pip install uv && uv sync --frozen --no-dev", TXT),
    ("COPY src/ src/", TXT),
    ("", TXT),
    ("# Stage 2: runtime (enxuto)", MUTED),
    ("FROM python:3.13-slim AS runtime", TXT),
    ("RUN useradd -r screenflix        # non-root", OK),
    ("COPY --from=builder /app/.venv .venv/", TXT),
    ("HEALTHCHECK CMD curl -f .../healthz", OK),
    ("CMD [\"uvicorn\", \"screenflix.main:app\", ...]", TXT),
], size=13)
card(s, Inches(7.95), Inches(2.15), Inches(4.75), Inches(4.4),
     "Por que multi-stage?",
     ["Ferramentas de build ficam fora",
      "da imagem final → menor e segura.",
      "",
      "• Ordem dos COPY = cache de camadas",
      "• Usuário non-root (segurança)",
      "• HEALTHCHECK em /healthz",
      "  (contrato de infra)",
      "• PYTHONPATH=/app/src"],
     accent=BRAND, body_size=13)

# 1.3 compose
s = content_slide("Dia 1 · Infra & Containers", "docker-compose: db + api",
                  "Dois serviços ativos — o api sobe após o banco ficar saudável")
code_box(s, Inches(0.6), Inches(2.35), Inches(7.4), Inches(4.2), [
    ("services:", ACCENT),
    ("  db:", TXT),
    ("    image: postgres:18-alpine", TXT),
    ("    volumes:", TXT),
    ("      - ./scripts/init.sql:/.../init.sql", TXT),
    ("    ports: [\"5532:5432\"]", TXT),
    ("    healthcheck: pg_isready ...", OK),
    ("  api:", TXT),
    ("    build: { dockerfile: Dockerfile }", TXT),
    ("    ports: [\"8000:8000\"]", TXT),
    ("    volumes: [ ./src:/app/src ]  # hot-reload", ACCENT),
    ("    depends_on:", TXT),
    ("      db: { condition: service_healthy }", OK),
], size=13)
card(s, Inches(8.25), Inches(2.35), Inches(4.45), Inches(2.0),
     "image vs build",
     ["db usa imagem pronta;", "api constrói do Dockerfile."], accent=ACCENT, body_size=13)
card(s, Inches(8.25), Inches(4.5), Inches(4.45), Inches(2.05),
     "YAGNI na prática",
     ["O serviço redis, antes reservado,", "foi removido por falta de uso.",
      "Não carregue infra que não usa."], accent=WARN, body_size=13)

# 1.4 comandos + makefile
s = content_slide("Dia 1 · Infra & Containers", "Comandos & Makefile")
code_box(s, Inches(0.6), Inches(2.15), Inches(6.0), Inches(4.4), [
    ("# Comandos do dia a dia", MUTED),
    ("docker build -t screenflix .", TXT),
    ("docker compose up --build", TXT),
    ("docker compose logs -f api", TXT),
    ("docker compose exec db psql ...", TXT),
    ("docker compose down", TXT),
    ("", TXT),
    ("# Com Makefile vira:", MUTED),
    ("make up      make check", OK),
    ("make dev     make build", OK),
], size=14)
card(s, Inches(6.85), Inches(2.15), Inches(5.85), Inches(4.4),
     "Makefile — interface única",
     ["Ninguém decora o comando longo do",
      "uvicorn. Com Makefile vira make dev.",
      "",
      "• make help → documentação executável",
      "• make check → roda os quality gates",
      "  (ruff format, ruff check, mypy, pytest)",
      "",
      "⚠ NÃO existe ainda no repo:",
      "   criá-lo é o hands-on do dia."],
     accent=BRAND, body_size=13)

# ══════════════════════════════════════════════════════════════
# DIA 2
# ══════════════════════════════════════════════════════════════
divider("Dia 2", "Programação Pragmática",
        ["Como o código se organiza por dentro",
         "Ambientes modulares · Clean Architecture · sem dogmatismo"], ACCENT)

# 2.1 ambientes modulares
s = content_slide("Dia 2 · Programação Pragmática", "Ambientes modulares", accent=ACCENT)
txt(s, Inches(0.6), Inches(2.1), Inches(12), Inches(0.5),
    [[R("A aplicação ", 16, MUTED), R("não sabe", 16, TXT, True),
      R(" onde roda — ela apenas lê configuração de fora.", 16, MUTED)]])
code_box(s, Inches(0.6), Inches(2.8), Inches(7.0), Inches(3.5), [
    ("class Settings(BaseSettings):", TXT),
    ("    model_config = SettingsConfigDict(", TXT),
    ("        env_prefix=\"SCREENFLIX_\",", ACCENT),
    ("        env_file=\".env\")", TXT),
    ("    app_name: str        # obrigatório", TXT),
    ("    database_url: str     # obrigatório", TXT),
    ("    openai_api_key: str   # obrigatório", TXT),
    ("    database_echo: bool = False  # default", MUTED),
], size=13)
card(s, Inches(7.85), Inches(2.8), Inches(4.85), Inches(3.5),
     "Fail fast + seguro",
     ["• Pydantic valida na inicialização",
      "  → config inválida falha na hora",
      "• Campos sem default = impossível",
      "  subir mal configurado",
      "• Segredos vêm do ambiente;",
      "  .env no .gitignore",
      "",
      "Frontend: mesma ideia via",
      "VITE_API_BASE (com modo demo)."],
     accent=ACCENT, body_size=13)

# 2.2 camadas
s = content_slide("Dia 2 · Programação Pragmática", "Clean Architecture — as 4 camadas", accent=ACCENT)
layers = [
    ("Presentation", "API HTTP (FastAPI) — endpoints finos, validação, DI", BRAND),
    ("Application", "Services, Use Cases e Schemas (DTOs) — orquestração", ACCENT),
    ("Infrastructure", "Repositórios — persistência (BaseRepository genérico)", OK),
    ("Domain", "Entidades Media e Episode — o coração do negócio", WARN),
]
y = Inches(2.3)
for i, (name, desc, col) in enumerate(layers):
    yy = y + i*Inches(0.92)
    indent = Inches(0.6) + Inches(i)*Inches(0.0)
    rect(s, Inches(0.9), yy, Inches(11.5), Inches(0.78), CARD, line=SURF2, line_w=Pt(1))
    rect(s, Inches(0.9), yy, Inches(0.12), Inches(0.78), col)
    txt(s, Inches(1.2), yy + Inches(0.11), Inches(3.0), Inches(0.5), [[R(name, 18, col, True)]])
    txt(s, Inches(4.2), yy + Inches(0.14), Inches(8.0), Inches(0.5), [[R(desc, 14, MUTED)]])
txt(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.5),
    [[R("Dependências apontam para dentro:  ", 15, TXT, True),
      R("presentation → application → domain", 15, ACCENT, True, MONO)]])

# 2.3 regra de dependência / padrões
s = content_slide("Dia 2 · Programação Pragmática", "Padrões & regra de ouro", accent=ACCENT)
code_box(s, Inches(0.6), Inches(2.15), Inches(7.0), Inches(3.2), [
    ("@router.get(\"/movies/top5\")", TXT),
    ("async def top5(service = Depends(...)):", TXT),
    ("    return await service.top_five(\"movie\")", ACCENT),
    ("", TXT),
    ("class MediaService:", TXT),
    ("    async def top_five(self, media_type):", TXT),
    ("        return await self.repository.media \\", TXT),
    ("            .list_all(order_by=\"rating\", ...)", TXT),
], size=13)
card(s, Inches(7.85), Inches(2.15), Inches(4.85), Inches(3.2),
     "Injeção de Dependência",
     ["A DI encadeia:",
      "sessão → service → endpoint.",
      "O response_model garante o",
      "contrato da saída."], accent=ACCENT, body_size=14)
rect(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.85), SURF2)
txt(s, Inches(0.85), Inches(5.78), Inches(11.6), Inches(0.5),
    [[R("REGRA DE OURO:  ", 15, WARN, True),
      R("nunca consulte o banco direto de um endpoint — sempre service → repository.", 15, TXT)]])

# 2.4 trade-offs
s = content_slide("Dia 2 · Programação Pragmática", "Onde o ScreenFlix escolheu pragmatismo", accent=ACCENT)
bullets(s, Inches(0.7), Inches(2.35), Inches(11.8), [
    ("Entidades = modelos ORM  ", "— simplicidade sobre portabilidade (não são POPOs puras)."),
    ("register fire-and-forget  ", "— dispara asyncio.create_task e responde na hora; sem fila durável, dívida consciente."),
    ("Top 5 no frontend  ", "— apesar de existirem endpoints /top5, a Home calcula no cliente."),
], size=17, gap=16, marker_col=ACCENT)
rect(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.1), CARD, line=SURF2, line_w=Pt(1))
txt(s, Inches(0.95), Inches(5.55), Inches(11.4), Inches(0.8),
    [[R("Clean Architecture é um espectro, não um checklist. ", 16, TXT, True),
      R("O valor está em camadas com responsabilidades claras — e em saber, conscientemente, onde relaxar.", 16, MUTED)]])

# ══════════════════════════════════════════════════════════════
# DIA 3
# ══════════════════════════════════════════════════════════════
divider("Dia 3", "IA no Desenvolvimento",
        ["Fazer a IA respeitar a infra (Dia 1) e a arquitetura (Dia 2)",
         "Claude Code · Skills · Subagents · Agent Teams"], OK)

# 3.1 insight
s = content_slide("Dia 3 · IA no Desenvolvimento", "O insight do ScreenFlix",
                  "O repo já foi construído pensando em IA — vamos traduzir esse contexto", accent=OK)
rect(s, Inches(0.6), Inches(2.35), Inches(5.9), Inches(0.55), SURF2)
txt(s, Inches(0.8), Inches(2.44), Inches(5.6), Inches(0.4), [[R("Contexto genérico do repo", 14, MUTED, True)]])
rect(s, Inches(6.9), Inches(2.35), Inches(5.8), Inches(0.55), SURF2)
txt(s, Inches(7.1), Inches(2.44), Inches(5.5), Inches(0.4), [[R("Artefato nativo do Claude Code", 14, OK, True)]])
maps = [
    ("AGENTS.md", "CLAUDE.md (memória de projeto)"),
    ("AI-CONFIG.json → quality_checks", "Skill /quality-gate + Hook"),
    ("agents/backend-development.md", "Subagent  backend-dev"),
    ("agents/qa-guidelines.md", "Subagent  qa-reviewer"),
    ("agents/security-check.md", "Subagent  security-check"),
    ("vários subagents juntos", "Agent Team"),
]
y = Inches(3.0)
for i, (a, b) in enumerate(maps):
    yy = y + i*Inches(0.6)
    rect(s, Inches(0.6), yy, Inches(5.9), Inches(0.5), CARD)
    txt(s, Inches(0.8), yy + Inches(0.11), Inches(5.6), Inches(0.35), [[R(a, 13, TXT, fnt=MONO)]])
    txt(s, Inches(6.52), yy + Inches(0.07), Inches(0.35), Inches(0.35), [[R("→", 16, OK, True)]])
    rect(s, Inches(6.9), yy, Inches(5.8), Inches(0.5), CARD)
    txt(s, Inches(7.1), yy + Inches(0.11), Inches(5.5), Inches(0.35), [[R(b, 13, OK, fnt=MONO)]])

# 3.2 Skills vs Subagents
s = content_slide("Dia 3 · IA no Desenvolvimento", "Skills vs Subagents", accent=OK)
card(s, Inches(0.6), Inches(2.2), Inches(5.95), Inches(4.3),
     "Skill  ·  .claude/skills/*/SKILL.md",
     ["Procedimento reutilizável na",
      "conversa principal, carregado",
      "sob demanda. Vira /comando.",
      "",
      "Ex.: /quality-gate roda",
      "ruff format · ruff check ·",
      "mypy src · pytest",
      "(os gates bloqueantes do projeto).",
      "",
      "Hook: torna o gate automático",
      "no evento Stop (settings.json)."],
     accent=OK, body_size=14)
card(s, Inches(6.75), Inches(2.2), Inches(5.95), Inches(4.3),
     "Subagent  ·  .claude/agents/*.md",
     ["Trabalhador isolado, com contexto",
      "e ferramentas próprios. O Claude",
      "principal delega pela description.",
      "",
      "Cada guia de /agents vira um",
      "especialista:",
      "  architect · backend-dev ·",
      "  database-dev · security-check ·",
      "  qa-reviewer",
      "",
      "Roda do zero e devolve resultado."],
     accent=BRAND, body_size=14)

# 3.3 Agent Teams
s = content_slide("Dia 3 · IA no Desenvolvimento", "Agent Teams",
                  "Orquestração multi-agente — experimental", accent=OK)
pill(s, Inches(0.6), Inches(2.1), "EXPERIMENTAL", color=WARN, w=Inches(2.0))
# comparison table
headers = ["Aspecto", "Subagents", "Agent Teams"]
rows = [
    ("Comunicação", "Reportam ao principal", "Teammates conversam entre si"),
    ("Coordenação", "O principal gerencia", "Lista de tarefas compartilhada"),
    ("Uso ideal", "Tarefas focadas", "Trabalho paralelo complexo"),
]
tx = Inches(0.6); ty = Inches(2.75); cw = [Inches(2.8), Inches(4.55), Inches(4.75)]
xx = tx
for j, h in enumerate(headers):
    rect(s, xx, ty, cw[j], Inches(0.55), BRAND if j == 0 else SURF2)
    txt(s, xx + Inches(0.15), ty + Inches(0.12), cw[j]-Inches(0.2), Inches(0.4), [[R(h, 14, TXT, True)]])
    xx += cw[j]
for i, row in enumerate(rows):
    xx = tx; yy = ty + Inches(0.55) + i*Inches(0.62)
    for j, cell in enumerate(row):
        rect(s, xx, yy, cw[j], Inches(0.62), CARD if i % 2 == 0 else SURF1)
        col = ACCENT if j == 0 else MUTED
        bold = j == 0
        txt(s, xx + Inches(0.15), yy + Inches(0.15), cw[j]-Inches(0.25), Inches(0.4), [[R(cell, 13, col, bold)]])
        xx += cw[j]
txt(s, Inches(0.6), Inches(5.7), Inches(12), Inches(0.6),
    [[R("Regra prática: ", 15, OK, True),
      R("comece simples — subagent para mudanças focadas; team só para trabalho amplo e paralelizável.", 15, MUTED)]])

# ══════════════════════════════════════════════════════════════
# FECHAMENTO
# ══════════════════════════════════════════════════════════════
s = content_slide("Síntese", "A jornada, num só repositório")
data = [
    ("DIA 1", "Empacotar a app", "Dockerfile, compose, Makefile", BRAND),
    ("DIA 2", "Organizar por dentro", "Ambientes modulares + Clean Architecture", ACCENT),
    ("DIA 3", "IA disciplinada", "CLAUDE.md, Skills, Subagents, Teams", OK),
]
y = Inches(2.4)
for i, (d, t, desc, col) in enumerate(data):
    yy = y + i*Inches(1.05)
    rect(s, Inches(0.7), yy, Inches(1.7), Inches(0.85), col)
    txt(s, Inches(0.7), yy + Inches(0.24), Inches(1.7), Inches(0.5), [[R(d, 18, TXT, True)]], align=PP_ALIGN.CENTER)
    rect(s, Inches(2.5), yy, Inches(10.1), Inches(0.85), CARD, line=SURF2, line_w=Pt(1))
    txt(s, Inches(2.75), yy + Inches(0.12), Inches(9.6), Inches(0.4), [[R(t, 17, TXT, True)]])
    txt(s, Inches(2.75), yy + Inches(0.46), Inches(9.6), Inches(0.35), [[R(desc, 13, MUTED)]])
txt(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.9),
    [[R("Os artefatos de IA fazem cumprir a arquitetura: ", 15, TXT, True),
      R("o qa-reviewer exige testes, o architect protege as camadas, a Skill roda os gates.", 15, MUTED)]])

# Slide final — obrigado
s = prs.slides.add_slide(BLANK)
bg(s, INK)
rect(s, 0, 0, SW, Inches(0.16), BRAND)
rect(s, 0, Inches(7.34), SW, Inches(0.16), BRAND)
txt(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.2), [[R("Obrigado!", 46, TXT, True)]])
txt(s, Inches(0.95), Inches(3.75), Inches(11.5), Inches(0.6),
    [[R("Perguntas? Vamos colocar a mão na massa.", 18, MUTED)]])
txt(s, Inches(0.95), Inches(4.6), Inches(11.5), Inches(0.5),
    [[R("Material completo:  ", 14, MUTED), R("docs/workshop/ (roteiros + apostila)", 14, ACCENT, True, MONO)]])

prs.save("/home/user/screenflix/scratchpad/ScreenFlix-Workshop.pptx")
print("OK — slides:", len(prs.slides._sldIdLst))
